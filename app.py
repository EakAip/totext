# 接口：8015

#文件，图片，音频，视频

# 转文字

import os
import torch
import PyPDF2
import argparse
import subprocess
import numpy as np
import gradio as gr
from tqdm import tqdm
from docx import Document
from funasr import AutoModel
from pptx import Presentation
from vision.seeit import draw_box
from vision import OCR, init_in_out
from moviepy.editor import VideoFileClip



def extract_text_from_docx(docx_file):
    doc = Document(docx_file)
    all_text = []
    for para in doc.paragraphs:
        all_text.append(para.text.replace('\n', ' '))
    return '\n'.join(all_text)

def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text.replace("\n", " "))
    return '\n'.join(all_text)

def extract_text_from_pdf(pdf_file):
    with open(pdf_file, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        all_text = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                all_text.append(text.replace('\n', ' '))
            else:
                all_text.append("No text found on this page.")
    return '\n'.join(all_text)

def convert_doc_to_docx(doc_file):
    doc_dir = os.path.dirname(doc_file)
    new_file_name = os.path.basename(doc_file).replace(".doc", ".docx")
    new_file_path = os.path.join(doc_dir, new_file_name)
    command = [
        "libreoffice", "--headless", "--convert-to", 
        "docx:Office Open XML Text", doc_file, "--outdir", doc_dir
    ]
    subprocess.run(command, check=True)
    return new_file_path

def extract_text_from_ocr(input_file, output_dir):
    args = argparse.Namespace(inputs=input_file, output_dir=output_dir)
    ocr = OCR()
    images, outputs = init_in_out(args)
    
    all_text = []

    for i, img in enumerate(tqdm(images, desc="Processing images")):
        if img.mode == 'RGBA':
            img = img.convert('RGB')
        bxs = ocr(np.array(img))
        bxs = [(line[0], line[1][0]) for line in bxs]
        bxs = [{
            "text": t,
            "bbox": [b[0][0], b[0][1], b[1][0], b[-1][1]],
            "type": "ocr",
            "score": 1} for b, t in bxs if b[0][0] <= b[1][0] and b[0][1] <= b[-1][1]]
        img = draw_box(img, bxs, ["ocr"], 1.)
        img.save(outputs[i], quality=95)
        # 输出img路径
        print(outputs[i])
        page_text = "\n".join([o["text"] for o in bxs])
        with open(outputs[i] + ".txt", "w+") as f:
            f.write(page_text)
        
        all_text.append(page_text)

    return "\n".join(all_text), outputs[-1]

def extract_text_from_audio(input_audio):
    model = AutoModel(
        model="damo/speech_seaco_paraformer_large_asr_nat-zh-cn-16k-common-vocab8404-pytorch",
        model_revision="v2.0.4", 
        vad_model="damo/speech_fsmn_vad_zh-cn-16k-common-pytorch",
        vad_model_revision="v2.0.4",
        punc_model="damo/punc_ct-transformer_zh-cn-common-vocab272727-pytorch",
        punc_model_revision="v2.0.4",
    )

    # 尝试释放未使用的GPU内存
    torch.cuda.empty_cache()

    res = model.generate(input=input_audio)
    return res[0]["text"]

def process_file(file_obj, file_type):
    print("File object received:", file_obj)
    output_dir = "./ocr_outputs"  # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)

    if file_type == "DOCX":
        return extract_text_from_docx(file_obj), None
    elif file_type == "DOC":
        docx_file_path = convert_doc_to_docx(file_obj)
        return extract_text_from_docx(docx_file_path), None
    elif file_type == "PPT":
        return extract_text_from_ppt(file_obj), None
    elif file_type == "PDF":
        return extract_text_from_pdf(file_obj), None
    elif file_type == "PDF-OCR":
        # 文件路径已经是完整路径，可以直接使用
        return extract_text_from_ocr(file_obj, output_dir)
    else:
        return "Unsupported file type.", None

def process_audio(audio_file):
    return extract_text_from_audio(audio_file)

def extract_audio_from_video(video_path):
    print("************************开始视频转音频************************")
    # 为提取的音频创建一个临时文件名，使用原视频路径的文件名部分
    audio_file_path = os.path.splitext(video_path)[0] + ".wav"
    
    # 使用VideoFileClip加载视频文件
    video_clip = VideoFileClip(video_path)
    
    # 获取视频的音频部分
    audio_clip = video_clip.audio
    
    # 将音频保存到指定的文件路径
    audio_clip.write_audiofile(audio_file_path, codec='pcm_s16le')
    
    # 释放资源
    audio_clip.close()
    video_clip.close()
    
    # 返回音频文件的路径
    return audio_file_path

def process_video(video_file):
    audio_path = extract_audio_from_video(video_file)
    return extract_text_from_audio(audio_path)
    
with gr.Blocks() as demo:
    gr.Markdown("<h1 style='text-align: center;'></h1>")
    gr.Markdown("<h1 style='text-align: center;'>🍸文件文本提取</h1>")
    gr.Markdown("<h1 style='text-align: center;'></h1>")
    gr.Markdown("<h1 style='text-align: center;'></h1>")
    gr.Markdown("## 选择文件类型")

    with gr.Tabs():
        with gr.TabItem("上传文档"):
            with gr.Row():
                with gr.Column():
                    file_input = gr.File(label="上传文件")
                    file_type = gr.Radio(choices=["DOCX", "DOC", "PPT", "PDF", "PDF-OCR"], label="选择文件类型", value="DOCX")
                    gr.Examples(["/home/root-jyd/jyd/totext/data/数据结构C语言版.docx",
                                 "/home/root-jyd/jyd/totext/data/数据结构实验指导书.doc",
                                 "/home/root-jyd/jyd/totext/data/数据结构第一章.pptx",
                                 "/home/root-jyd/jyd/totext/data/Python核心编程第三版.pdf",
                                 "/home/root-jyd/jyd/totext/data/数据库系统概论-扫描件.pdf"
                                 ], [file_input],label='示例文档')
                    
                    submit_file_button = gr.Button("提取文本")
                with gr.Column(scale=2):
                    output_file_text = gr.Textbox(label="提取的文本", lines=20)
                    output_file_image = gr.Image(label="处理后的图片", type="filepath", visible=False)
            
        with gr.TabItem("上传图片"):
            with gr.Row():
                with gr.Column():
                    image_input = gr.Image(label="上传图片", type="filepath")
                    gr.Examples(["/home/root-jyd/jyd/totext/data/测试1.png"
                                 ], [image_input],label='示例图片demo')
                    submit_image_button = gr.Button("提取文本")
                with gr.Column():
                    output_image = gr.Image(label="处理", type="filepath")
                    output_image_text = gr.Textbox(label="提取的文本", lines=20)
                    
        
        with gr.TabItem("上传音频"):
            with gr.Row():
                with gr.Column():
                    audio_input = gr.Audio(label="上传音频", type="filepath")
                    gr.Examples(["/home/root-jyd/jyd/totext/data/小对话.wav"
                                 ], [audio_input],label='示例音频demo')
                    submit_audio_button = gr.Button("提取文本")
                with gr.Column():
                    output_audio_text = gr.Textbox(label="提取的文本", lines=20)
        
        with gr.TabItem("上传视频"):
            with gr.Row():
                with gr.Column():
                    video_input = gr.Video(label="上传视频")
                    gr.Examples([
                                "/home/root-jyd/jyd/totext/data/录屏.mp4"
                                ],
                                [video_input], label='示例视频demo')
                    submit_video_button = gr.Button("提取文本")
                with gr.Column():
                    output_video_text = gr.Textbox(label="提取的文本", lines=20)

    def handle_file(file_obj, file_type):
        if file_obj is not None:
            text, image_path = process_file(file_obj, file_type)
            if image_path:
                return text, image_path
            else:
                return text, None
        else:
            return "请上传文件。", None

    def handle_image(image_obj):
        if image_obj is not None:
            return process_file(image_obj, "PDF-OCR")
        else:
            return "请上传图片。", None

    def handle_audio(audio_obj):
        if audio_obj is not None:
            return process_audio(audio_obj)
        else:
            return "请上传音频文件。"
    
    def handle_video(video_obj):
        if video_obj is not None:
            return process_video(video_obj)
        else:
            return "请上传视频文件。"

    submit_file_button.click(
        fn=handle_file,
        inputs=[file_input, file_type],
        outputs=[output_file_text, output_file_image]
    )
    
    submit_image_button.click(
        fn=handle_image,
        inputs=image_input,
        outputs=[output_image_text, output_image]
    )
    
    submit_audio_button.click(
        fn=handle_audio,
        inputs=audio_input,
        outputs=output_audio_text
    )
    
    submit_video_button.click(
        fn=handle_video,
        inputs=video_input,
        outputs=output_video_text
    )

demo.launch(server_port=8015, server_name='0.0.0.0',show_error=True)
