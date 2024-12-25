from pptx import Presentation

def extract_text_from_ppt(ppt_file):
    # 创建 Presentation 对象
    prs = Presentation(ppt_file)
    
    # 存储所有文本的列表
    all_text = []
    
    # 遍历每个幻灯片
    for slide in prs.slides:
        # 遍历每个形状（包括文本框、表格等）
        for shape in slide.shapes:
            # 检查形状是否包含文本
            if hasattr(shape, "text"):
                # 如果是文本框，将文本添加到列表中
                all_text.append(shape.text.replace("\n", " "))
    
    return all_text

# 指定你的 PPT 文件路径
ppt_file_path = "/home/jyd01/wangruihua/knowledge_ner/data/数据结构第一章.pptx"

# 调用函数并打印结果
text_from_ppt = extract_text_from_ppt(ppt_file_path)
# 保存为txt文件
with open('/home/jyd01/wangruihua/knowledge_ner/data/数据结构第一章.txt', 'w', encoding='utf-8') as f:
    for text in text_from_ppt:
        f.write(text)
    f.close()

print(text_from_ppt)
